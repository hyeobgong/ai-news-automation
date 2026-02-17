from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_text_format(shape, font_size, is_bold=False, color=RGBColor(255, 255, 255)):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = font_size
        paragraph.font.bold = is_bold
        paragraph.font.color.rgb = color

# Dark Theme Colors
BG_COLOR = RGBColor(30, 30, 30)  # Dark Gray
TEXT_COLOR = RGBColor(255, 255, 255)  # White
TITLE_SIZE = Pt(32)
BODY_SIZE = Pt(18)


def create_presentation():
    prs = Presentation()

    # --- 1. Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI 뉴스 자동화 프로젝트 기획안"
    subtitle.text = "매일 아침 9시, 나만의 AI 비서 만들기\n(Slack 연동 & Gemini 기반)"
    
    set_text_format(title, Pt(40), True, TEXT_COLOR)
    set_text_format(subtitle, Pt(24), False, TEXT_COLOR)

    # --- 2. Project Rules (Revised) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "1. 프로젝트 정의 및 규칙 (Rules)"
    content.text = (
        "1) 알림 채널: Slack (슬랙)\n"
        "   - 업무용 메신저 활용, Webhook 연동 용이\n"
        "2) 콘텐츠 포맷: [제목] + [3줄 요약] + [원문 링크]\n"
        "   - 가독성 최우선, 이미지 프리뷰 포함\n"
        "3) 언어 설정: 한국어 요약 기본 (영어 원문 옵션)\n"
        "   - Gemini 프롬프트로 '한국어 번역 요약' 지시\n"
        "4) 저작권 준수: 전문(Full-text) 수집 금지\n"
        "   - 요약문과 링크만 제공하여 저작권 이슈 회피"
    )
    set_text_format(title, TITLE_SIZE, True, TEXT_COLOR)
    set_text_format(content, BODY_SIZE, False, TEXT_COLOR)

    # --- 3. Tech Stack ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "2. 단계별 추천 도구 (Tech Stack)"
    content.text = (
        "① 수집 (Collector): Python (feedparser)\n"
        "   - Google News RSS 등에서 24시간 이내 뉴스 필터링\n\n"
        "② 지능형 처리 (AI Brain): Google Gemini API (Flash)\n"
        "   - 역할: 번역, 3줄 요약, 중복 기사 통합\n\n"
        "③ 자동화 (Scheduler): GitHub Actions\n"
        "   - 매일 오전 9시(KST) 자동 실행 (서버 비용 0원)\n\n"
        "④ 알림 (Notifier): Slack Incoming Webhook\n"
        "   - 텔레그램 대체, 깔끔한 UI 제공"
    )
    set_text_format(title, TITLE_SIZE, True, TEXT_COLOR)
    set_text_format(content, BODY_SIZE, False, TEXT_COLOR)

    # --- 4. Workflow Diagram (Text Based) ---
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    title.text = "3. 전체 워크플로우 (Process)"
    set_text_format(title, TITLE_SIZE, True, TEXT_COLOR)

    # Draw Boxes
    shapes = slide.shapes

    def style_shape(shape, fill_color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        # Black text for light backgrounds
        set_text_format(shape, Pt(14), True, RGBColor(0, 0, 0))

    # Box 1
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(2)
    height = Inches(1)
    box1 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box1.text = "09:00 AM Trigger\n(GitHub Actions)"
    style_shape(box1, RGBColor(176, 196, 222))  # Light Steel Blue

    # Arrow 1
    arrow1 = shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(2.9), Inches(0.8), Inches(0.2)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = TEXT_COLOR

    # Box 2
    box2 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), top, width, height)
    box2.text = "뉴스 수집 & 필터링\n(Python)"
    style_shape(box2, RGBColor(152, 251, 152))  # Pale Green

    # Arrow 2
    arrow2 = shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(2.9), Inches(0.8), Inches(0.2)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = TEXT_COLOR

    # Box 3
    box3 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), top, width, height)
    box3.text = "AI 요약 & 번역\n(Gemini API)"
    style_shape(box3, RGBColor(135, 206, 250))  # Light Sky Blue

    # Arrow 3 (Down)
    arrow3 = shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(8), Inches(3.6), Inches(0.2), Inches(0.8)
    )
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = TEXT_COLOR

    # Box 4
    box4 = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4.5), width, height
    )
    box4.text = "Slack 알림 전송\n(Webhook)"
    style_shape(box4, RGBColor(255, 182, 193))  # Light Pink

    # --- 5. Mind Map (Visual) ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    title.text = "4. 프로젝트 마인드맵 (Mind Map)"
    set_text_format(title, TITLE_SIZE, True, TEXT_COLOR)

    shapes = slide.shapes

    # Center Node
    center_x = Inches(5)
    center_y = Inches(3.5)
    center_box = shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x - Inches(1.25),
        center_y - Inches(0.75),
        Inches(2.5),
        Inches(1.5),
    )
    center_box.text = "AI 뉴스\n자동화"
    style_shape(center_box, RGBColor(255, 215, 0)) # Gold

    def add_white_connector(from_shape, to_shape, from_idx, to_idx):
        conn = shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            from_shape.left + (from_shape.width * from_idx[0]),
            from_shape.top + (from_shape.height * from_idx[1]),
            to_shape.left + (to_shape.width * to_idx[0]),
            to_shape.top + (to_shape.height * to_idx[1]),
        )
        conn.line.color.rgb = TEXT_COLOR
        return conn

    # Branch 1: 수집 (Top Left)
    branch1 = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x - Inches(3.5),
        center_y - Inches(2),
        Inches(2),
        Inches(1),
    )
    branch1.text = "수집 (Source)\n- Google RSS\n- TechCrunch"
    style_shape(branch1, RGBColor(176, 196, 222))
    
    # Custom connector logic replacing the hardcoded coordinates
    connector1 = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        center_box.left,
        center_box.top + Inches(0.75), # Center left of oval roughly
        branch1.left + Inches(2),
        branch1.top + Inches(0.5), # Center right of rect
    )
    connector1.line.color.rgb = TEXT_COLOR

    # Branch 2: 처리 (Top Right)
    branch2 = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x + Inches(1.5),
        center_y - Inches(2),
        Inches(2),
        Inches(1),
    )
    branch2.text = "처리 (AI)\n- Gemini Flash\n- 한국어 요약\n- 중복 제거"
    style_shape(branch2, RGBColor(144, 238, 144)) # Light Green

    connector2 = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        center_box.left + Inches(2.5),
        center_box.top + Inches(0.75),
        branch2.left,
        branch2.top + Inches(0.5),
    )
    connector2.line.color.rgb = TEXT_COLOR

    # Branch 3: 알림 (Bottom Left)
    branch3 = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x - Inches(3.5),
        center_y + Inches(1),
        Inches(2),
        Inches(1),
    )
    branch3.text = "알림 (Output)\n- Slack\n- 제목+3줄요약\n- 원문 링크"
    style_shape(branch3, RGBColor(255, 182, 193)) # Light Pink
    
    connector3 = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        center_box.left,
        center_box.top + Inches(0.75), # Reuse center y anchor
        branch3.left + Inches(2),
        branch3.top + Inches(0.5),
    )
    connector3.line.color.rgb = TEXT_COLOR

    # Branch 4: 인프라 (Bottom Right)
    branch4 = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x + Inches(1.5),
        center_y + Inches(1),
        Inches(2),
        Inches(1),
    )
    branch4.text = "인프라 (Infra)\n- Python\n- GitHub Actions\n- 09:00 KST"
    style_shape(branch4, RGBColor(221, 160, 221)) # Plum
    
    connector4 = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        center_box.left + Inches(2.5),
        center_box.top + Inches(0.75),
        branch4.left,
        branch4.top + Inches(0.5),
    )
    connector4.line.color.rgb = TEXT_COLOR

    # --- 6. Example Output ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "5. 최종 결과물 예시 (Slack)"
    content.text = (
        "🤖 [AI Daily Briefing] 2026-02-17\n\n"
        "1. OpenAI, 차세대 추론 모델 'o3' 발표\n"
        "   - 기존 o1 모델 대비 수학적 추론 능력 30% 향상\n"
        "   - 기업용 API 가격 50% 인하\n"
        "   👉 [원문 보기 (TechCrunch)](https://example.com)\n\n"
        "2. 구글, Gemini 2.0 업데이트\n"
        "   - 멀티모달 기능 강화로 비디오 인식 속도 2배 증가\n"
        "   - 개발자용 1M 컨텍스트 윈도우 무료 제공\n"
        "   👉 [원문 보기 (Google Blog)](https://example.com)"
    )
    set_text_format(title, TITLE_SIZE, True, TEXT_COLOR)
    set_text_format(content, BODY_SIZE, False, TEXT_COLOR)

    prs.save("AI_News_Automation_Plan_Dark.pptx")
    print("PPT 파일이 생성되었습니다: AI_News_Automation_Plan_Dark.pptx")


if __name__ == "__main__":
    create_presentation()
