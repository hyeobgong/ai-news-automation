from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_progress_report():
    prs = Presentation()

    # --- Style Constants ---
    BG_COLOR = RGBColor(30, 30, 30)       # Dark Background
    TEXT_COLOR = RGBColor(240, 240, 240)  # Light Text
    ACCENT_COLOR = RGBColor(100, 180, 255) # Light Blue Accent
    TITLE_FONT_SIZE = Pt(36)
    BODY_FONT_SIZE = Pt(20)

    def apply_dark_theme(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def style_text_frame(text_frame, font_size=BODY_FONT_SIZE, is_bold=False):
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.color.rgb = TEXT_COLOR
            paragraph.font.bold = is_bold

    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        content.text = content_text
        style_text_frame(content.text_frame)

    # --- 1. Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_theme(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI 뉴스 자동화 프로젝트\n진행 상황 보고서"
    subtitle.text = "2026년 2월 18일\n작성자: AI 뉴스 자동화 팀"
    
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    style_text_frame(subtitle.text_frame, Pt(24))

    # --- 2. Project Overview ---
    add_slide(
        "1. 프로젝트 개요",
        "목표:\n"
        "- 매일 쏟아지는 AI 관련 뉴스를 자동 수집 및 요약\n"
        "- 사용자에게 핵심 정보만 간추려 제공 (시간 절약)\n\n"
        "핵심 기능:\n"
        "- Google News RSS 기반 실시간 뉴스 수집\n"
        "- TF-IDF 알고리즘을 통한 중복 기사 제거\n"
        "- Gemini Pro/Flash API를 활용한 3줄 요약 및 중요도 평가"
    )

    # --- 3. Development Progress ---
    add_slide(
        "2. 개발 진행 상황 (Phase 1 & 2 완료)",
        "✅ 기획 및 문서화 (Completed)\n"
        "   - 프로젝트 계획 PPT 생성기 (make_ppt.py) 구현\n"
        "   - GitHub 리포지토리 설정 및 README 작성\n\n"
        "✅ 뉴스 수집기 (Collector)\n"
        "   - '생성형 AI' 등 키워드 기반 RSS 크롤링 구현 (feedparser)\n\n"
        "✅ 뉴스 처리기 (Processor)\n"
        "   - 중복 기사 필터링 (TF-IDF + Cosine Similarity)\n"
        "   - Gemini API 연동을 통한 기사 분류 및 요약 로직 구현"
    )

    # --- 4. Current Issues & Solutions ---
    add_slide(
        "3. 주요 이슈 및 해결 방안",
        "🚨 이슈 1: Gemini API 할당량 초과 (429 Error)\n"
        "   - 원인: 무료 티어(Free Tier) 사용량 제한 도달\n"
        "   - 해결: API 오류 발생 시 프로그램이 멈추지 않도록 예외 처리 강화\n"
        "   - 대안: 더미 데이터(Dummy Data) 생성 로직을 추가하여 전체 파이프라인 테스트 가능하도록 조치\n\n"
        "⚠️ 이슈 2: 중복 제거 정확도\n"
        "   - 현황: 제목이 유사해도 키워드 차이로 다른 기사로 인식\n"
        "   - 계획: 유사도 임계값(Threshold) 조정 (0.6 → 0.4) 및 비교 대상 확대"
    )

    # --- 5. Future Plans ---
    add_slide(
        "4. 향후 계획 (Next Steps)",
        "1️⃣ API 안정화\n"
        "   - 유료 API 키 전환 또는 새 키 발급을 통한 모델 정상화\n\n"
        "2️⃣ 데이터 저장소 구축\n"
        "   - 분석된 뉴스를 로컬 파일(JSON) 또는 DB에 저장하는 구조 마련\n\n"
        "3️⃣ 리포트 생성 및 배포 자동화\n"
        "   - 매일 아침 자동 실행되도록 GitHub Actions 워크플로우 설정\n"
        "   - 이메일 또는 슬랙(Slack)으로 요약본 전송 기능 추가"
    )

    # --- 6. Tech Stack ---
    add_slide(
        "5. 사용 기술 스택",
        "💻 언어: Python 3.13\n"
        "📊 문서 자동화: python-pptx (PPT 생성)\n"
        "🌐 수집: feedparser (RSS 크롤링)\n"
        "🧠 AI 모델: Google Gemini (google-genai) 2.0/1.5 Flash\n"
        "🧮 분석: scikit-learn (TF-IDF, Cosine Similarity)\n"
        "⚙️ 관리: Git, GitHub, python-dotenv"
    )

    # Save
    output_file = "AI_News_Progress_Report.pptx"
    prs.save(output_file)
    print(f"보고서 생성 완료: {output_file}")

if __name__ == "__main__":
    create_progress_report()
