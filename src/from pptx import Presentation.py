from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()

    # --- 1. Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI 뉴스 자동화 프로젝트 기획안"
    subtitle.text = "매일 아침 9시, 나만의 AI 비서 만들기\n(Slack 연동 & Gemini 기반)"

    # --- 2. Project Rules (Revised) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "1. 프로젝트 정의 및 규칙 (Rules)"
    content.text = (
        "1) 알림 채널: Slack (슬랙)\n"
        "   - 업무용 메신저 활용, Webhook 연동 용이\n"
        "2) 콘텐츠 포맷: [제목] + [3줄 요약] + [원문 링크]\n"
        "   - 가독성 최우선, 이미지 프리뷰 포함\n"
        "3) 언어 설정: 한국어 요약 기본 (영어 원문 옵션)\n"
        "   - Gemini 프롬프트로 '한국어 번역 요약' 지시\n"
        "4) 저작권 준수: 전문(Full-text) 수집 금지\n"
        "   - 요약문과 링크만 제공하여 저작권 이슈 회피"
    )

    # --- 3. Tech Stack ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "2. 단계별 추천 도구 (Tech Stack)"
    content.text = (
        "① 수집 (Collector): Python (feedparser)\n"
        "   - Google News RSS 등에서 24시간 이내 뉴스 필터링\n\n"
        "② 지능형 처리 (AI Brain): Google Gemini API (Flash)\n"
        "   - 역할: 번역, 3줄 요약, 중복 기사 통합\n\n"
        "③ 자동화 (Scheduler): GitHub Actions\n"
        "   - 매일 오전 9시(KST) 자동 실행 (서버 비용 0원)\n\n"
        "④ 알림 (Notifier): Slack Incoming Webhook\n"
        "   - 텔레그램 대체, 깔끔한 UI 제공"
    )

    # --- 4. Workflow Diagram (Text Based) ---
    slide_layout = prs.slide_layouts[5] # Blank
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "3. 전체 워크플로우 (Process)"

    # Draw Boxes
    shapes = slide.shapes
    
    # Box 1
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(2)
    height = Inches(1)
    box1 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box1.text = "09:00 AM Trigger\n(GitHub Actions)"

    # Arrow 1
    arrow1 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(2.9), Inches(0.8), Inches(0.2))

    # Box 2
    box2 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), top, width, height)
    box2.text = "뉴스 수집 & 필터링\n(Python)"

    # Arrow 2
    arrow2 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(2.9), Inches(0.8), Inches(0.2))

    # Box 3
    box3 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), top, width, height)
    box3.text = "AI 요약 & 번역\n(Gemini API)"
    
    # Arrow 3 (Down)
    arrow3 = shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(8), Inches(3.6), Inches(0.2), Inches(0.8))

    # Box 4
    box4 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4.5), width, height)
    box4.text = "Slack 알림 전송\n(Webhook)"

    # --- 5. Mind Map (Visual) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "4. 프로젝트 마인드맵 (Mind Map)"

    shapes = slide.shapes
    
    # Center Node
    center_x = Inches(5)
    center_y = Inches(3.5)
    center_box = shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(1.25), center_y - Inches(0.75), Inches(2.5), Inches(1.5))
    center_box.text = "AI 뉴스\n자동화"
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = RGBColor(255, 204, 0) # Gold

    # Branch 1: 수집 (Top Left)
    branch1 = shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(3.5), center_y - Inches(2), Inches(2), Inches(1))
    branch1.text = "수집 (Source)\n- Google RSS\n- TechCrunch"
    branch1.fill.solid()
    branch1.fill.fore_color.rgb = RGBColor(173, 216, 230) # Light Blue
    shapes.add_connector(MSO_SHAPE.STRAIGHT_CONNECTOR_1, center_box.left, center_box.top + Inches(0.5), branch1.left + Inches(2), branch1.top + Inches(0.5))

    # Branch 2: 처리 (Top Right)
    branch2 = shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x + Inches(1.5), center_y - Inches(2), Inches(2), Inches(1))
    branch2.text = "처리 (AI)\n- Gemini Flash\n- 한국어 요약\n- 중복 제거"
    branch2.fill.solid()
    branch2.fill.fore_color.rgb = RGBColor(144, 238, 144) # Light Green
    shapes.add_connector(MSO_SHAPE.STRAIGHT_CONNECTOR_1, center_box.left + Inches(2.5), center_box.top + Inches(0.5), branch2.left, branch2.top + Inches(0.5))

    # Branch 3: 알림 (Bottom Left)
    branch3 = shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(3.5), center_y + Inches(1), Inches(2), Inches(1))
    branch3.text = "알림 (Output)\n- Slack\n- 제목+3줄요약\n- 원문 링크"
    branch3.fill.solid()
    branch3.fill.fore_color.rgb = RGBColor(255, 182, 193) # Light Pink
    shapes.add_connector(MSO_SHAPE.STRAIGHT_CONNECTOR_1, center_box.left, center_box.top + Inches(1), branch3.left + Inches(2), branch3.top + Inches(0.5))

    # Branch 4: 인프라 (Bottom Right)
    branch4 = shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x + Inches(1.5), center_y + Inches(1), Inches(2), Inches(1))
    branch4.text = "인프라 (Infra)\n- Python\n- GitHub Actions\n- 09:00 KST"
    branch4.fill.solid()
    branch4.fill.fore_color.rgb = RGBColor(221, 160, 221) # Plum
    shapes.add_connector(MSO_SHAPE.STRAIGHT_CONNECTOR_1, center_box.left + Inches(2.5), center_box.top + Inches(1), branch4.left, branch4.top + Inches(0.5))

    # --- 6. Example Output ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "5. 최종 결과물 예시 (Slack)"
    content.text = (
        "🤖 [AI Daily Briefing] 2026-02-17\n\n"
        "1. OpenAI, 차세대 추론 모델 'o3' 발표\n"
        "   - 기존 o1 모델 대비 수학적 추론 능력 30% 향상\n"
        "   - 기업용 API 가격 50% 인하\n"
        "   👉 [원문 보기 (TechCrunch)](https://example.com)\n\n"
        "2. 구글, Gemini 2.0 업데이트\n"
        "   - 멀티모달 기능 강화로 비디오 인식 속도 2배 증가\n"
        "   - 개발자용 1M 컨텍스트 윈도우 무료 제공\n"
        "   👉 [원문 보기 (Google Blog)](https://example.com)"
    )

    prs.save('AI_News_Automation_Plan.pptx')
    print("PPT 파일이 생성되었습니다: AI_News_Automation_Plan.pptx")

if __name__ == "__main__":
    create_presentation()